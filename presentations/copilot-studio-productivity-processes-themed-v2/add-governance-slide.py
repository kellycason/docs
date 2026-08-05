from __future__ import annotations

import argparse
import os
import shutil
from copy import deepcopy
from datetime import datetime
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DECK_PATH = Path(__file__).with_name(
    "Copilot Studio - From Productivity to Processes - Themed v2.pptx"
)
ASSET_DIR = Path(__file__).with_name("_assets")
ARCHIVE_DIR = Path(__file__).with_name("_archive")
SLIDE_TITLE = "Govern agents with confidence"


COLORS = {
    "ink": "091F2C",
    "bg": "F2F1F7",
    "white": "FFFFFF",
    "blue": "0078D4",
    "teal": "008272",
    "sky": "8DC8E8",
    "purple": "8661C5",
    "magenta": "C03BC4",
    "coral": "FF5C39",
    "muted": "5B6672",
    "line": "DDDCE8",
    "shadow": "DAD8E5",
    "pale_blue": "EAF4FB",
    "pale_purple": "ECEAF6",
    "footer": "9AA3AD",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float,
    color: str = "ink",
    bold: bool = False,
    font: str = "Segoe UI",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_rounded_rect(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str,
    line: str,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.color.rgb = rgb(COLORS[line])
    shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = "line",
    width: float = 1.6,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(COLORS[color])
    connector.line.width = Pt(width)
    return connector


def add_pillar_card(
    slide,
    *,
    x: float,
    y: float,
    width: float,
    title: str,
    glyph: str,
    value: str,
    bullets: list[str],
    accent: str,
):
    height = 1.62
    add_rounded_rect(
        slide,
        x + 0.035,
        y + 0.045,
        width,
        height,
        fill="shadow",
        line="shadow",
    )
    card = add_rounded_rect(
        slide,
        x,
        y,
        width,
        height,
        fill="white",
        line="line",
        line_width=0.8,
    )
    card.name = f"Pillar - {title}"

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.13),
        Inches(y),
        Inches(width - 0.26),
        Inches(0.055),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(COLORS[accent])
    accent_bar.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.18),
        Inches(y + 0.18),
        Inches(0.38),
        Inches(0.38),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(COLORS[accent])
    badge.line.fill.background()
    add_text(
        slide,
        glyph,
        x + 0.18,
        y + 0.18,
        0.38,
        0.38,
        size=14.5,
        color="white",
        font="Segoe MDL2 Assets",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title.upper(),
        x + 0.67,
        y + 0.19,
        width - 0.87,
        0.3,
        size=12.2,
        color="ink",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        value,
        x + 0.18,
        y + 0.58,
        width - 0.36,
        0.34,
        size=9.6,
        color="ink",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )

    column_gap = 0.1
    column_width = (width - 0.4 - column_gap) / 2
    for index, bullet in enumerate(bullets):
        column = index % 2
        row = index // 2
        bullet_x = x + 0.2 + column * (column_width + column_gap)
        bullet_y = y + 1.0 + row * 0.27
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(bullet_x),
            Inches(bullet_y + 0.087),
            Inches(0.06),
            Inches(0.06),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(COLORS[accent])
        dot.line.fill.background()
        add_text(
            slide,
            bullet,
            bullet_x + 0.13,
            bullet_y,
            column_width - 0.13,
            0.23,
            size=8.5,
            color="muted",
            valign=MSO_ANCHOR.MIDDLE,
        )


def find_power_platform_logo(prs: Presentation) -> bytes | None:
    source_slide = prs.slides[2]
    for shape in source_slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            and shape.top > Inches(5.8)
            and shape.left < Inches(2.0)
        ):
            return shape.image.blob
    return None


def set_speaker_notes(prs: Presentation, slide, notes: str):
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        source_notes_slide = next(
            (
                source_slide.notes_slide
                for source_slide in prs.slides
                if source_slide is not slide
                and source_slide.notes_slide.notes_text_frame is not None
            ),
            None,
        )
        if source_notes_slide is None:
            raise RuntimeError("No existing notes placeholder is available to clone.")
        for shape in source_notes_slide.shapes:
            if shape.is_placeholder:
                notes_slide.shapes._spTree.insert_element_before(
                    deepcopy(shape._element), "p:extLst"
                )
    if notes_slide.notes_text_frame is None:
        raise RuntimeError("Unable to create a notes body placeholder.")
    notes_slide.notes_text_frame.text = notes


def add_foundation_item(
    slide,
    *,
    x: float,
    label: str,
    glyph: str,
    accent: str,
):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(6.39),
        Inches(0.27),
        Inches(0.27),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(COLORS[accent])
    badge.line.fill.background()
    add_text(
        slide,
        glyph,
        x,
        6.39,
        0.27,
        0.27,
        size=10.3,
        color="white",
        font="Segoe MDL2 Assets",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        label,
        x + 0.36,
        6.39,
        1.05,
        0.27,
        size=9.0,
        color="muted",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )


def is_governance_slide(slide) -> bool:
    return any(
        SLIDE_TITLE.lower() in shape.text.lower()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def remove_governance_slides(prs: Presentation) -> int:
    indices = [
        index for index, slide in enumerate(prs.slides) if is_governance_slide(slide)
    ]
    for index in reversed(indices):
        slide_id = prs.slides._sldIdLst[index]
        relationship_id = slide_id.rId
        prs.slides._sldIdLst.remove(slide_id)
        prs.part.drop_rel(relationship_id)
    return len(indices)


def append_governance_slide(prs: Presentation):
    if any(is_governance_slide(slide) for slide in prs.slides):
        raise RuntimeError("Governance slide already exists; no changes were made.")

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.name = "Govern agents with confidence"

    add_text(
        slide,
        "ENTERPRISE GOVERNANCE",
        0.85,
        0.48,
        4.2,
        0.25,
        size=10.8,
        color="blue",
        bold=True,
        font="Segoe UI Semibold",
    )
    add_text(
        slide,
        SLIDE_TITLE,
        0.85,
        0.82,
        8.8,
        0.48,
        size=26.5,
        color="ink",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "One platform for security, governance, compliance, lifecycle management, and operational visibility.",
        0.85,
        1.38,
        11.75,
        0.34,
        size=12.4,
        color="muted",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )

    halo = add_rounded_rect(
        slide,
        4.35,
        2.08,
        4.63,
        1.65,
        fill="pale_purple",
        line="pale_purple",
    )
    halo.name = "Copilot Studio governance halo"

    add_line(slide, 4.18, 2.73, 4.55, 2.91, color="blue", width=1.5)
    add_line(slide, 4.18, 4.55, 4.55, 3.3, color="purple", width=1.5)
    add_line(slide, 9.15, 2.73, 8.78, 2.91, color="magenta", width=1.5)
    add_line(slide, 9.15, 4.55, 8.78, 3.3, color="coral", width=1.5)
    add_line(slide, 6.665, 3.67, 6.665, 3.78, color="teal", width=1.8)

    hub = add_rounded_rect(
        slide,
        4.55,
        2.35,
        4.23,
        1.32,
        fill="ink",
        line="blue",
        line_width=1.2,
    )
    hub.name = "Copilot Studio governance hub"

    chip = add_rounded_rect(
        slide,
        5.84,
        2.13,
        1.65,
        0.35,
        fill="purple",
        line="purple",
    )
    chip.name = "Govern label"
    add_text(
        slide,
        "GOVERN",
        5.84,
        2.13,
        1.65,
        0.35,
        size=9.4,
        color="white",
        bold=True,
        font="Segoe UI Semibold",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    slide.shapes.add_picture(
        str(ASSET_DIR / "studio.png"),
        Inches(4.93),
        Inches(2.68),
        width=Inches(0.67),
        height=Inches(0.67),
    )
    add_text(
        slide,
        "COPILOT STUDIO",
        5.83,
        2.64,
        2.55,
        0.4,
        size=18.5,
        color="white",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Secure by design. Governed at scale.",
        5.83,
        3.08,
        2.55,
        0.3,
        size=10.2,
        color="sky",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_pillar_card(
        slide,
        x=0.75,
        y=1.92,
        width=3.43,
        title="Identity",
        glyph=chr(0xE716),
        value="Control who can build, publish, and use agents",
        bullets=["Entra ID", "Authentication", "RBAC", "Environment security"],
        accent="blue",
    )
    add_pillar_card(
        slide,
        x=0.75,
        y=3.78,
        width=3.43,
        title="Security",
        glyph=chr(0xE72E),
        value="Control what agents can access",
        bullets=[
            "DLP Policies",
            "Connector governance",
            "Environment isolation",
            "Data boundaries",
        ],
        accent="purple",
    )
    add_pillar_card(
        slide,
        x=9.15,
        y=1.92,
        width=3.43,
        title="Compliance",
        glyph=chr(0xE8FB),
        value="Audit and investigate AI interactions",
        bullets=["Microsoft Purview", "Audit logs", "eDiscovery", "DLP & DSPM"],
        accent="magenta",
    )
    add_pillar_card(
        slide,
        x=9.15,
        y=3.78,
        width=3.43,
        title="Operations",
        glyph=chr(0xE713),
        value="Measure usage, cost, and performance",
        bullets=[
            "Analytics",
            "Consumption monitoring",
            "Agent inventory",
            "Runtime insights",
        ],
        accent="coral",
    )
    add_pillar_card(
        slide,
        x=4.55,
        y=3.78,
        width=4.23,
        title="Lifecycle",
        glyph=chr(0xEC7A),
        value="Manage agents like enterprise applications",
        bullets=["Solutions", "Dev/Test/Prod", "ALM", "Managed deployment"],
        accent="teal",
    )

    callout = add_rounded_rect(
        slide,
        0.75,
        5.56,
        12.06,
        0.48,
        fill="pale_blue",
        line="pale_blue",
    )
    callout.name = "Money statement"
    add_rounded_rect(
        slide,
        0.92,
        5.635,
        1.78,
        0.33,
        fill="ink",
        line="ink",
    )
    add_text(
        slide,
        "ONE PLATFORM",
        0.92,
        5.635,
        1.78,
        0.33,
        size=9.0,
        color="white",
        bold=True,
        font="Segoe UI Semibold",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Move from experimentation to production with the same platform.",
        2.96,
        5.61,
        9.36,
        0.37,
        size=13.1,
        color="ink",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )

    banner = add_rounded_rect(
        slide,
        0.75,
        6.2,
        11.83,
        0.61,
        fill="white",
        line="line",
        line_width=0.8,
    )
    banner.name = "Power Platform foundation"
    for x, width, color in [
        (0.75, 3.943, "blue"),
        (4.693, 3.943, "purple"),
        (8.636, 3.944, "magenta"),
    ]:
        segment = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(6.2),
            Inches(width),
            Inches(0.055),
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = rgb(COLORS[color])
        segment.line.fill.background()

    power_platform_logo = find_power_platform_logo(prs)
    if power_platform_logo:
        slide.shapes.add_picture(
            BytesIO(power_platform_logo),
            Inches(1.02),
            Inches(6.31),
            width=Inches(0.38),
            height=Inches(0.38),
        )
    add_text(
        slide,
        "Built on Microsoft Power Platform",
        1.55,
        6.31,
        2.72,
        0.35,
        size=10.7,
        color="ink",
        bold=True,
        font="Segoe UI Semibold",
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_foundation_item(slide, x=4.46, label="Identity", glyph=chr(0xE716), accent="blue")
    add_foundation_item(slide, x=5.93, label="Security", glyph=chr(0xE72E), accent="blue")
    add_foundation_item(slide, x=7.42, label="Governance", glyph=chr(0xE713), accent="purple")
    add_foundation_item(slide, x=9.25, label="ALM", glyph=chr(0xEC7A), accent="magenta")
    add_foundation_item(slide, x=10.43, label="Analytics", glyph=chr(0xE721), accent="coral")

    add_text(
        slide,
        "Copilot Studio",
        0.85,
        7.12,
        2.0,
        0.17,
        size=8.5,
        color="footer",
    )
    add_text(
        slide,
        "11",
        12.28,
        7.08,
        0.35,
        0.18,
        size=9,
        color="footer",
        align=PP_ALIGN.RIGHT,
    )

    notes = (
        "Primary talk track:\n"
        '"What makes Copilot Studio different is that governance isn\'t bolted on later. '
        "It's built into the platform from day one. Administrators can control who builds agents, "
        "what data they access, how they're deployed, how they're monitored, and how AI interactions "
        "are audited. That allows organizations to move from experimentation to production with confidence.\"\n\n"
        "Pillar prompts:\n"
        "Identity - Not everyone should be able to build, publish, or administer agents. Copilot Studio "
        "inherits enterprise identity and access controls through Power Platform.\n"
        "Security - Administrators decide what data, tools, connectors, and actions an agent can access.\n"
        "Lifecycle - Agents move through governed environments just like any other enterprise solution.\n"
        "Operations - Organizations need visibility into adoption, activity, and consumption before they can scale AI.\n"
        "Compliance - Purview provides the compliance layer for AI, enabling auditability, investigations, "
        "data protection, and compliance workflows.\n\n"
        "Close: You don't need another governance platform. Governance is built in.\n\n"
        "Sources:\n"
        "https://learn.microsoft.com/en-us/microsoft-copilot-studio/security-and-governance\n"
        "https://learn.microsoft.com/en-us/purview/ai-copilot-studio\n"
        "https://learn.microsoft.com/en-us/microsoft-copilot-studio/admin-logging-copilot-studio\n"
        "https://microsoft.sharepoint.com/teams/BSWHExternal/Shared%20Documents/Microsoft%20BSWH%20AI%20SWAT%20Team/"
        "Copilot%20Studio%20Agent%20Governance%20&%20Security%20Controls%20TDM%20(L300).pdf?web=1"
    )
    set_speaker_notes(prs, slide, notes)
    return slide


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--replace",
        action="store_true",
        help="Replace an existing governance slide instead of refusing to duplicate it.",
    )
    args = parser.parse_args()

    if not DECK_PATH.exists():
        raise FileNotFoundError(DECK_PATH)

    prs = Presentation(DECK_PATH)
    original_slide_count = len(prs.slides)
    removed_slide_count = 0
    if args.replace:
        removed_slide_count = remove_governance_slides(prs)
    append_governance_slide(prs)
    expected_slide_count = original_slide_count - removed_slide_count + 1

    ARCHIVE_DIR.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup_label = "pre-governance-polish" if args.replace else "pre-governance"
    backup_path = ARCHIVE_DIR / f"{backup_label}-{timestamp}.pptx"
    shutil.copy2(DECK_PATH, backup_path)

    temp_path = DECK_PATH.with_name(f"{DECK_PATH.stem}.governance.tmp.pptx")
    prs.save(temp_path)

    verification = Presentation(temp_path)
    if len(verification.slides) != expected_slide_count:
        raise RuntimeError("Saved deck does not contain exactly one governance slide.")
    final_slide_text = "\n".join(
        shape.text
        for shape in verification.slides[-1].shapes
        if getattr(shape, "has_text_frame", False)
    )
    if SLIDE_TITLE not in final_slide_text:
        raise RuntimeError("Saved deck is missing the governance slide title.")

    os.replace(temp_path, DECK_PATH)
    print(f"Updated: {DECK_PATH}")
    print(f"Backup:  {backup_path}")
    print(f"Slides:  {len(verification.slides)}")


if __name__ == "__main__":
    main()